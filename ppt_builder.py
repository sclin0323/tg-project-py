"""
ppt_builder.py — v2: 從 data_loader 取得 report 並產出 PPT
適配新版 WBS schema, 並重新設計第二頁為四區塊任務細項視圖
"""
from __future__ import annotations
from datetime import date, timedelta
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_AUTO_SIZE

from ppt_components import (
    C, HEALTH, STATUS_COLOR, FONT, SLIDE_W, SLIDE_H,
    rect, round_rect, circle, text, divider,
    page_header, kpi_card, progress_bar, status_pill, health_badge,
    fmt_date, fmt_date_short,
)
from data_loader import HEALTH_LABELS, STATUS_DONE, STATUS_RISK


def _blank_slide(pres):
    return pres.slides.add_slide(pres.slide_layouts[6])


# ══════════════════════════════════════════════════════
# Slide 1: 封面
# ══════════════════════════════════════════════════════
def slide_cover(pres, r):
    s = _blank_slide(pres)
    overall = r["overall"]
    main_c  = HEALTH[overall]["main"]
    soft_c  = HEALTH[overall]["soft"]

    rect(s, 0, 0, 0.10, SLIDE_H, main_c)

    text(s, 0.55, 1.2, 6.0, 0.9, "專案管理週報",
         size=42, bold=True, color=C["text"])
    text(s, 0.55, 2.05, 6.0, 0.34, "Executive Status Report",
         size=13, italic=True, color=C["muted"])

    rect(s, 0.55, 2.62, 2.8, 0.008, C["divider"])

    round_rect(s, 0.55, 2.80, 4.0, 0.70, soft_c,
               line=main_c, line_w=1.2, radius=0.18)
    text(s, 0.55, 2.80, 4.0, 0.70, r["headline"],
         size=14, bold=True, color=main_c,
         align="center", valign="middle")

    text(s, 0.55, 3.78, 5.0, 0.34,
         f"報告日期　{fmt_date(r['report_date'])}",
         size=12, color=C["muted"])
    text(s, 0.55, 4.16, 5.0, 0.34,
         f"下次回報　{fmt_date(r['next_report'])}",
         size=12, color=C["muted"])

    rect(s, 7.0, 0, SLIDE_W - 7.0, SLIDE_H, C["bg"])
    rect(s, 7.0, 0, SLIDE_W - 7.0, 0.06, C["divider"])

    items = [
        ("健康", r["stats"]["green_projects"], C["green"]),
        ("注意", r["stats"]["amber_projects"], C["amber"]),
        ("警報", r["stats"]["red_projects"],   C["red"]),
    ]
    for i, (label, n, col) in enumerate(items):
        ty = 1.3 + i * 1.15
        circle(s, 7.45, ty + 0.18, 0.30, col)
        text(s, 7.75, ty, 1.6, 0.45, str(n),
             size=32, bold=True, color=col, valign="middle")
        text(s, 7.75, ty + 0.50, 1.8, 0.26, f"{label} 專案",
             size=10, color=C["muted"])


# ══════════════════════════════════════════════════════
# Slide 2: 執行摘要
# ══════════════════════════════════════════════════════
def slide_executive_summary(pres, r):
    s = _blank_slide(pres)
    page_header(s, "執行摘要",
                f"專案一覽　·　{fmt_date(r['report_date'])}",
                accent=C["text"])

    cards = [
        ("需立即關注", r["stats"]["red_projects"],   C["red"],   C["redSoft"]),
        ("需要觀察",   r["stats"]["amber_projects"], C["amber"], C["amberSoft"]),
        ("進度健康",   r["stats"]["green_projects"], C["green"], C["greenSoft"]),
    ]
    for i, (label, n, col, soft) in enumerate(cards):
        x = 0.3 + i * 3.15
        rect(s, x, 0.92, 2.95, 0.85, soft, line=C["divider"], line_w=0.5)
        rect(s, x, 0.92, 0.07, 0.85, col)
        text(s, x + 0.2, 0.92, 0.9, 0.85, str(n),
             size=36, bold=True, color=col, valign="middle")
        text(s, x + 1.1, 0.92, 1.8, 0.85, f"{label}　專案",
             size=11, color=C["muted"], valign="middle")

    TX, TY, TW = 0.3, 1.98, 9.4
    COLS = [0.4, 1.85, 0.85, 1.55, 1.30, 0.90, 2.55]
    HDRS = ["", "專案 / PM", "完成率", "進度 / 時間", "預估完工", "偏差", "本期關鍵"]

    hx = TX
    for h, w in zip(HDRS, COLS):
        rect(s, hx, TY, w, 0.32, C["bg"], line=C["divider"], line_w=0.4)
        text(s, hx, TY, w, 0.32, h,
             size=9, bold=True, color=C["muted"], align="center", valign="middle")
        hx += w

    ROW_H = 0.55
    for i, p in enumerate(r["projects"]):
        y = TY + 0.32 + i * ROW_H
        bg = C["white"] if i % 2 == 0 else C["rowAlt"]
        rect(s, TX, y, TW, ROW_H, bg, line=C["divider"], line_w=0.3)

        h = p["health"]
        hc = HEALTH[h["level"]]["main"]

        rx = TX
        circle(s, rx + COLS[0]/2, y + ROW_H/2, 0.22, hc)
        rx += COLS[0]

        text(s, rx + 0.08, y + 0.06, COLS[1] - 0.16, 0.27, p["name"],
             size=11, bold=True, color=C["text"], valign="middle")
        text(s, rx + 0.08, y + 0.30, COLS[1] - 0.16, 0.22,
             f"{p['owner']}  ·  {h['label']}",
             size=8.5, color=hc)
        rx += COLS[1]

        text(s, rx, y, COLS[2], ROW_H, f"{round(h['done_pct'])}%",
             size=16, bold=True, color=C["text"],
             align="center", valign="middle")
        rx += COLS[2]

        bx, bw = rx + 0.12, COLS[3] - 0.24
        by = y + ROW_H/2 - 0.05
        progress_bar(s, bx, by, bw, 0.13, h["done_pct"], hc,
                     vs_pct=h["elapsed_pct"], vs_color=C["text"])
        text(s, bx, y + ROW_H - 0.18, bw, 0.16,
             f"時間 {round(h['elapsed_pct'])}%",
             size=7, color=C["light"], align="center")
        rx += COLS[3]

        if h.get("forecast_available", False):
            fc_txt = fmt_date(h["forecast_end"])
            fc_col = C["text"]
        else:
            fc_txt = "—"
            fc_col = C["light"]
        text(s, rx, y + 0.06, COLS[4], 0.27, fc_txt,
             size=10, bold=True, color=fc_col,
             align="center", valign="middle")
        text(s, rx, y + 0.30, COLS[4], 0.22,
             f"目標 {fmt_date(p['target'])}",
             size=7.5, color=C["light"], align="center")
        rx += COLS[4]

        if not h.get("forecast_available", False):
            txt, col = "資料不足", C["light"]
        else:
            d = h["forecast_drift"]
            if d > 0:   txt, col = f"+{d} 天", C["red"]
            elif d < 0: txt, col = f"{d} 天", C["green"]
            else:       txt, col = "準時", C["green"]
        text(s, rx, y, COLS[5], ROW_H, txt,
             size=12 if "資料" not in txt else 10,
             bold=True, color=col, align="center", valign="middle")
        rx += COLS[5]

        if p["decisions"]:
            top_d = p["decisions"][0]
            key_txt = top_d["issue"]
            key_col = C["red"] if top_d["severity"] == "高" else C["amber"]
            italic = False
        else:
            key_txt = "進度正常"
            key_col = C["green"]
            italic = True
        text(s, rx + 0.1, y, COLS[6] - 0.2, ROW_H, key_txt,
             size=9, color=key_col, italic=italic, valign="middle")

    bot_y = TY + 0.32 + len(r["projects"]) * ROW_H + 0.15
    text(s, 0.3, bot_y, 9.4, 0.26,
         f"需主管裁示事項：{r['stats']['decision_count']} 件　·　詳見各專案 Page 2 與彙整頁",
         size=9, color=C["light"], italic=True)


# ══════════════════════════════════════════════════════
# Slide 3..N: 每專案 Page 1 — 整體概況
# ══════════════════════════════════════════════════════
def slide_project_overview(pres, p, r):
    s = _blank_slide(pres)
    h  = p["health"]
    hc = HEALTH[h["level"]]["main"]
    soft = HEALTH[h["level"]]["soft"]
    deep = HEALTH[h["level"]]["deep"]

    period = f"{fmt_date(p['start'])} — {fmt_date(p['target'])}"
    page_header(s, p["name"],
                f"PM　{p['owner']}　·　{period}",
                page_tag="專案概況　1 / 2", accent=hc)

    LX, LW = 0.25, 4.40
    CY     = 0.85

    HERO_H   = 1.55
    METRIC_Y = CY + HERO_H + 0.08
    METRIC_H = 0.90
    PROG_Y   = METRIC_Y + METRIC_H + 0.08
    PROG_H   = 0.36
    DIST_Y   = PROG_Y + PROG_H + 0.10
    DIST_H   = SLIDE_H - DIST_Y - 0.30

    # Hero
    rect(s, LX, CY, LW, HERO_H, C["white"], line=C["divider"], line_w=0.5)
    health_badge(s, LX + 0.42, CY + 0.36, 0.34, h["level"], with_ring=True)
    text(s, LX + 0.70, CY + 0.14, 2.5, 0.46, h["label"],
         size=24, bold=True, color=hc, valign="middle")
    text(s, LX + 0.70, CY + 0.54, LW - 0.85, 0.24, p["stage"] or "—",
         size=10, color=C["muted"], valign="middle")

    sub_x = LX + 0.18
    sub_y = CY + 0.86
    sub_w = LW - 0.36
    sub_h = HERO_H - 0.86 - 0.10
    round_rect(s, sub_x, sub_y, sub_w, sub_h, soft,
               line=hc, line_w=0.6, radius=0.12)
    text(s, sub_x + 0.15, sub_y, sub_w - 0.3, sub_h, p["summary"],
         size=11, bold=True, color=deep, valign="middle")

    # 三大 KPI
    metrics = []
    metrics.append({"label": "完成率", "value": f"{round(h['done_pct'])}%", "color": hc})
    if h.get("forecast_available", False):
        metrics.append({"label": "預估完工",
                        "value": fmt_date_short(h["forecast_end"]),
                        "color": C["text"]})
    else:
        metrics.append({"label": "預估完工（資料不足）",
                        "value": "—", "color": C["light"]})
    if not h.get("forecast_available", False):
        drift_v, drift_c = "—", C["light"]
    else:
        d = h["forecast_drift"]
        if d > 0:   drift_v, drift_c = f"+{d} 天", C["red"]
        elif d < 0: drift_v, drift_c = f"{d} 天", C["green"]
        else:       drift_v, drift_c = "準時", C["green"]
    metrics.append({"label": "與目標偏差", "value": drift_v, "color": drift_c})

    mw = LW / 3
    for i, m in enumerate(metrics):
        kpi_card(s, LX + i * mw, METRIC_Y, mw, METRIC_H,
                 m["value"], m["label"],
                 value_color=m["color"], accent_color=None,
                 value_size=22, label_size=9.5)

    # 進度 vs 時間
    text(s, LX, PROG_Y, 2.0, 0.18, "進度 vs 時間",
         size=8.5, bold=True, color=C["muted"])
    text(s, LX + 2.0, PROG_Y, LW - 2.0, 0.18,
         f"完成 {round(h['done_pct'])}%　·　時間 {round(h['elapsed_pct'])}%",
         size=8.5, color=C["muted"], align="right")
    progress_bar(s, LX, PROG_Y + 0.20, LW, 0.16,
                 h["done_pct"], hc,
                 vs_pct=h["elapsed_pct"], vs_color=C["text"])

    # 任務分佈 (新版有 5 種狀態:完成/進行中/延遲/風險/待開始)
    rect(s, LX, DIST_Y, LW, DIST_H, C["white"],
         line=C["divider"], line_w=0.5)
    text(s, LX + 0.18, DIST_Y + 0.06, LW - 0.36, 0.22, "任務分佈",
         size=9.5, bold=True, color=C["muted"], valign="middle")
    text(s, LX + 0.18, DIST_Y + 0.06, LW - 0.36, 0.22,
         f"共 {p['total']} 個任務",
         size=8.5, color=C["light"], align="right", valign="middle")

    bar_x = LX + 0.20
    bar_y = DIST_Y + 0.40
    bar_w = LW - 0.40
    bar_h = 0.20
    total = max(1, p["total"])
    segs = [
        ("完成",   p["done"],        C["green"]),
        ("進行中", p["in_progress"], C["blue"]),
        ("延遲",   p["delayed"],     C["red"]),
        ("風險",   p["risk_count"],  C["amber"]),
        ("待開始", p["not_started"], C["divider"]),
    ]
    cx = bar_x
    for _, val, col in segs:
        w = bar_w * val / total
        if w > 0.01:
            rect(s, cx, bar_y, w, bar_h, col)
        cx += w

    # 圖例 (五項並排)
    legend_y = bar_y + bar_h + 0.12
    lx = bar_x
    seg_w = bar_w / len(segs)
    for label, val, col in segs:
        rect(s, lx, legend_y + 0.04, 0.12, 0.12, col)
        text(s, lx + 0.16, legend_y, seg_w - 0.18, 0.22,
             f"{label} {val}",
             size=8, color=C["muted"], valign="middle")
        lx += seg_w

    # 右半部: 里程碑甘特
    GX, GY = 4.80, CY
    GW = SLIDE_W - GX - 0.25
    GH = SLIDE_H - GY - 0.30
    _draw_gantt(s, GX, GY, GW, GH, p["milestones"], hc, r["report_date"])


def _draw_gantt(s, x, y, w, h, milestones, accent, today):
    """階段甘特圖 (新版)
       每一列 = 一個第 1 層階段, 橫條從 m['start'] 畫到 m['end'],
       依 progress (0~100) 決定填色比例:
         - 完成段 (左側): 用狀態色實心填滿
         - 未完成段 (右側): 淡色填底
       狀態色:
         - 完成   → 綠
         - 進行中 → 藍
         - 延遲   → 紅
         - 風險   → 琥珀
         - 待開始 → 灰
    """
    rect(s, x, y, w, h, C["white"], line=C["divider"], line_w=0.5)
    text(s, x + 0.2, y + 0.10, 3.0, 0.30, "階段進度",
         size=12, bold=True, color=C["text"], valign="middle")
    divider(s, x, y + 0.46, w)

    if not milestones:
        text(s, x, y + h/2 - 0.2, w, 0.4, "（無階段資料）",
             size=11, italic=True, color=C["light"], align="center")
        return

    # 時間軸範圍: 取所有階段日期 + 今日, 確保今日線在範圍內
    all_dates = [d for m in milestones for d in (m["start"], m["end"])]
    all_dates.append(today)
    minD, maxD = min(all_dates), max(all_dates)
    # 兩端 buffer, 避免貼邊
    span_days = max(1, (maxD - minD).days)
    pad = max(3, int(span_days * 0.02))
    minD = minD - timedelta(days=pad)
    maxD = maxD + timedelta(days=pad)
    total_d = max(1, (maxD - minD).days)

    LABEL_W = 2.10   # 階段名稱比較長, 多給一點寬度避免換行
    AREA_X  = x + LABEL_W
    AREA_W  = w - LABEL_W - 0.20
    HDR_Y   = y + 0.50

    def date_to_x(d):
        return AREA_X + (d - minD).days / total_d * AREA_W

    # 月份標尺
    cur = date(minD.year, minD.month, 1)
    while cur <= maxD:
        lx = date_to_x(cur)
        if AREA_X - 0.02 <= lx <= AREA_X + AREA_W + 0.02:
            rect(s, lx, HDR_Y, 0.006, h - 0.5, C["divider"])
            text(s, lx + 0.02, HDR_Y + 0.02, 0.4, 0.20,
                 f"{cur.month}月",
                 size=7, color=C["light"])
        if cur.month == 12:
            cur = date(cur.year + 1, 1, 1)
        else:
            cur = date(cur.year, cur.month + 1, 1)

    # 今日座標, 線延後到所有橫條後面再畫
    today_x = None
    if minD <= today <= maxD:
        today_x = date_to_x(today)

    # 最多顯示 10 列 (有些專案會有 10+ 個階段, 例如 KH 升級案有 11 個)
    MAX_ROWS = 10
    rows = milestones[:MAX_ROWS]
    extra_count = len(milestones) - len(rows)
    # 為「另有 N 個」字樣保留底部空間
    reserve_bottom = 0.20 if extra_count > 0 else 0.0

    bar_area_h = h - 0.78 - reserve_bottom
    row_h = min(0.42, bar_area_h / max(len(rows), 1))
    row_y0 = y + 0.74

    # 淡色底色 (未完成段) — 用更柔和的灰避免跟分隔線混淆
    soft_bg = C["divider"]

    for i, m in enumerate(rows):
        ry = row_y0 + i * row_h
        if i % 2 == 1:
            rect(s, x, ry, w, row_h, C["rowAlt"])

        status = m["status"]
        col = STATUS_COLOR.get(status, C["blue"])
        progress = m.get("progress", 0)

        # 階段名稱 (硬性截斷, 避免換行擠壓橫條)
        nm = m["name"]
        if len(nm) > 9:
            nm = nm[:9] + "…"
        # WBS 編號用固定寬度欄位 (左對齊), 名稱跟在後面 - 確保多列名稱起點對齊
        text(s, x + 0.10, ry, 0.36, row_h, f"{m['wbs_no']}.",
             size=8, bold=True, color=C["muted"], align="right", valign="middle")
        name_box = text(s, x + 0.50, ry, LABEL_W - 1.15, row_h, nm,
             size=8, bold=True, color=C["text"], valign="middle")
        # 關閉文字框 auto_size, 避免長名稱撐大文字框覆蓋編號欄
        name_box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        name_box.text_frame.word_wrap = False
        # 進度百分比 (右側緊靠橫條起點)
        text(s, x + LABEL_W - 0.67, ry, 0.60, row_h, f"{progress}%",
             size=8, bold=True, color=col,
             align="right", valign="middle")

        # ── 橫條: 從 start 畫到 end ──
        x1 = date_to_x(m["start"])
        x2 = date_to_x(m["end"])
        bar_w = max(0.06, x2 - x1)
        bar_h_px = min(0.18, row_h * 0.52)
        bar_y = ry + (row_h - bar_h_px) / 2

        # 1. 先畫底層: 整段淡色背景 (代表階段範圍)
        round_rect(s, x1, bar_y, bar_w, bar_h_px, soft_bg, radius=0.4)

        # 2. 再依 progress 畫已完成部分 (左側)
        if progress > 0:
            done_w = bar_w * progress / 100
            # 太短的填色看不到, 至少給個最小寬度
            done_w = max(0.04, done_w)
            done_w = min(done_w, bar_w)  # 不超出橫條
            round_rect(s, x1, bar_y, done_w, bar_h_px, col, radius=0.4)

    # 最後畫今日線, 確保疊在所有橫條上方
    if today_x is not None:
        rect(s, today_x - 0.012, HDR_Y, 0.024, h - 0.60 - reserve_bottom, accent)
        # 「今日」標籤放在底部, 加白底避免被橫條遮住
        label_y = y + h - 0.32 - reserve_bottom
        rect(s, today_x - 0.20, label_y, 0.40, 0.20, C["white"])
        text(s, today_x - 0.30, label_y, 0.60, 0.20, "今日",
             size=7.5, bold=True, color=accent, align="center", valign="middle")

    # 底部「另有 N 個階段未顯示」提示
    if extra_count > 0:
        text(s, x + 0.18, y + h - 0.22, w - 0.36, 0.18,
             f"另有 {extra_count} 個階段未顯示, 詳見 WBS 明細",
             size=7.5, italic=True, color=C["light"], align="right",
             valign="middle")


# ══════════════════════════════════════════════════════
# Slide 4..N: 每專案 Page 2 — 任務細項 (新版四區塊)
# ══════════════════════════════════════════════════════
def slide_project_tasks(pres, p, r):
    """第二頁:四區塊任務細項
       左上 — 已延遲(紅), 左下 — 風險(琥珀), 右上 — 進行中(藍), 右下 — 近期完成(綠)
    """
    s = _blank_slide(pres)
    h  = p["health"]
    hc = HEALTH[h["level"]]["main"]
    page_header(s, p["name"],
                f"任務細項　·　{fmt_date(r['report_date'])}",
                page_tag="任務細項　2 / 2", accent=hc)

    act = p["activity"]

    # 頂部 KPI 條 (薄條型,只佔 0.42 吋)
    kpi_y, kpi_h = 0.84, 0.42
    items = [
        ("已延遲", len(act["delayed"]), C["red"]),
        ("風險",   len(act["risk"]),    C["amber"]),
        ("進行中", len(act["active"]),  C["blue"]),
        ("近 14 天完成", len(act["recent"]), C["green"]),
    ]
    kw = (SLIDE_W - 0.5) / len(items)
    for i, (lab, val, col) in enumerate(items):
        x = 0.25 + i * kw
        # 自繪薄條 (取代 kpi_card)
        rect(s, x, kpi_y, kw - 0.12, kpi_h, C["white"],
             line=C["divider"], line_w=0.5)
        rect(s, x, kpi_y, kw - 0.12, 0.05, col)
        # 大數字置左
        text(s, x + 0.20, kpi_y, 0.80, kpi_h, str(val),
             size=22, bold=True, color=col, valign="middle")
        # 標籤置中偏右
        text(s, x + 1.00, kpi_y, kw - 0.12 - 1.10, kpi_h, lab,
             size=10, color=C["muted"], valign="middle")

    # 四區塊 (2x2)
    PY = kpi_y + kpi_h + 0.08
    PH_total = SLIDE_H - PY - 0.24
    PH = (PH_total - 0.10) / 2

    LX = 0.25
    LW_total = SLIDE_W - 0.50
    LW = (LW_total - 0.10) / 2     # 兩欄寬度
    RX = LX + LW + 0.10

    # 左上: 已延遲 ─────────────────────
    _draw_task_panel(
        s, LX, PY, LW, PH,
        title=f"已延遲（{len(act['delayed'])}）",
        subtitle="已逾期且未完成,需立即處理",
        tasks=act["delayed"],
        title_color=C["red"],
        accent=C["red"],
        bg_soft=C["redSoft"],
        kind="delayed",
        today=r["report_date"],
        empty_msg="✓ 目前無延遲任務",
        empty_color=C["green"],
        empty_bg=C["greenSoft"],
    )

    # 左下: 風險 ────────────────────────
    _draw_task_panel(
        s, LX, PY + PH + 0.10, LW, PH,
        title=f"風險（{len(act['risk'])}）",
        subtitle="尚未發生但可能影響進度的任務",
        tasks=act["risk"],
        title_color=C["amber"],
        accent=C["amber"],
        bg_soft=C["amberSoft"],
        kind="risk",
        today=r["report_date"],
        empty_msg="✓ 目前無風險任務",
        empty_color=C["green"],
        empty_bg=C["greenSoft"],
    )

    # 右上: 進行中 ──────────────────────
    _draw_task_panel(
        s, RX, PY, LW, PH,
        title=f"進行中（{len(act['active'])}）",
        subtitle="目前正在執行的任務",
        tasks=act["active"],
        title_color=C["blue"],
        accent=C["blue"],
        bg_soft=C["blueSoft"],
        kind="active",
        today=r["report_date"],
        empty_msg="目前無進行中任務",
        empty_color=C["light"],
        empty_bg=C["bg"],
    )

    # 右下: 近期完成 ────────────────────
    _draw_task_panel(
        s, RX, PY + PH + 0.10, LW, PH,
        title=f"近 14 天完成（{len(act['recent'])}）",
        subtitle="本期已實際完成的任務",
        tasks=act["recent"],
        title_color=C["green"],
        accent=C["green"],
        bg_soft=C["greenSoft"],
        kind="recent",
        today=r["report_date"],
        empty_msg="本期無完成任務",
        empty_color=C["light"],
        empty_bg=C["bg"],
    )


def _draw_task_panel(s, x, y, w, h, *,
                     title, subtitle, tasks,
                     title_color, accent, bg_soft, kind, today,
                     empty_msg, empty_color, empty_bg):
    """四區塊用的通用面板。
       kind: 'delayed' | 'risk' | 'active' | 'recent' — 決定底部顯示哪種日期資訊
    """
    # 面板外框
    rect(s, x, y, w, h, C["white"], line=C["divider"], line_w=0.5)
    rect(s, x, y, w, 0.05, accent)

    # 標題列 (合併標題+副標到一行,省高度)
    text(s, x + 0.18, y + 0.08, w * 0.55, 0.26, title,
         size=11, bold=True, color=title_color, valign="middle")
    text(s, x + w * 0.55, y + 0.08, w * 0.43, 0.26, subtitle,
         size=8, italic=True, color=C["muted"], align="right", valign="middle")
    divider(s, x + 0.18, y + 0.38, w - 0.36)

    # 內容區
    body_y = y + 0.44
    body_h = h - 0.44 - 0.08

    if not tasks:
        # 空狀態
        box_h = 0.60
        round_rect(s, x + 0.30, body_y + (body_h - box_h)/2,
                   w - 0.60, box_h, empty_bg,
                   line=C["divider"], line_w=0.5, radius=0.12)
        text(s, x + 0.30, body_y + (body_h - box_h)/2,
             w - 0.60, box_h, empty_msg,
             size=11, bold=True, color=empty_color,
             align="center", valign="middle")
        return

    # 動態決定能顯示幾筆 — 在可用空間內顯示越多越好
    # 卡片高度有上下限:不會太擠也不會空蕩
    MIN_ITEM_H = 0.36   # 最緊湊高度
    IDEAL_ITEM_H = 0.50  # 理想卡高 (內容剛好不擠)

    n_total = len(tasks)
    extra_h_reserve = 0.22

    cap_full = max(1, int(body_h / MIN_ITEM_H))

    if n_total <= cap_full:
        shown_n = n_total
        extra_h = 0.0
    else:
        shown_n = max(1, int((body_h - extra_h_reserve) / MIN_ITEM_H))
        extra_h = extra_h_reserve

    # 卡片高度: 用理想高度,但若 shown_n × IDEAL 超過可用空間,就壓縮
    item_h = min(IDEAL_ITEM_H, (body_h - extra_h) / shown_n)
    shown = tasks[:shown_n]

    for i, t in enumerate(shown):
        iy = body_y + i * item_h
        ih = item_h - 0.04

        # 背景 (奇數列淺色)
        if i % 2 == 1:
            rect(s, x + 0.12, iy, w - 0.24, ih, C["rowAlt"])

        # 左側狀態色條
        rect(s, x + 0.18, iy + 0.04, 0.04, ih - 0.08, accent)

        content_x = x + 0.28
        content_w = w - 0.46

        # 第一行: WBS 編號 + 工作項目 (同一行,WBS 用色,項目粗體)
        wbs_no_w = 0.62
        text(s, content_x, iy + 0.02, wbs_no_w, 0.22,
             t["wbs_no"],
             size=8, bold=True, color=accent, valign="middle")
        nm = t["name"].lstrip("★ ").strip()
        text(s, content_x + wbs_no_w + 0.02, iy + 0.02,
             content_w - wbs_no_w - 0.05, 0.22, nm,
             size=9.5, bold=True, color=C["text"], valign="middle")

        # 第二行: 依 kind 顯示不同的資訊
        info_y = iy + 0.20
        info_h = ih - 0.20
        owner = t["owner"] or "未指派"

        if kind == "delayed":
            plan = fmt_date_short(t["plan_end"]) if t["plan_end"] else "—"
            days_late = (today - t["plan_end"]).days if t["plan_end"] and t["plan_end"] < today else 0
            line = f"{owner}　·　預定 {plan}　·　已逾期 {days_late} 天"
            text(s, content_x, info_y, content_w, info_h, line,
                 size=8, color=C["redDeep"], valign="middle")
        elif kind == "risk":
            ps = fmt_date_short(t["plan_start"]) if t["plan_start"] else "—"
            pe = fmt_date_short(t["plan_end"]) if t["plan_end"] else "—"
            note = t["note"]
            if note:
                note = note.replace("【風險】", "").replace("【延遲】", "").strip()
                if len(note) > 28:
                    note = note[:28] + "…"
                line = f"{owner}　·　{ps}~{pe}　·　{note}"
            else:
                line = f"{owner}　·　預定 {ps} ~ {pe}"
            text(s, content_x, info_y, content_w, info_h, line,
                 size=8, color=C["amberDeep"], valign="middle")
        elif kind == "active":
            pe = fmt_date_short(t["plan_end"]) if t["plan_end"] else "—"
            if t["plan_end"]:
                d_to_end = (t["plan_end"] - today).days
                if d_to_end >= 0:
                    tail = f"還有 {d_to_end} 天"
                else:
                    tail = f"已逾期 {-d_to_end} 天"
            else:
                tail = ""
            line = f"{owner}　·　預定截止 {pe}　·　{tail}" if tail else f"{owner}　·　預定截止 {pe}"
            text(s, content_x, info_y, content_w, info_h, line,
                 size=8, color=C["muted"], valign="middle")
        else:  # recent
            ad = fmt_date_short(t["actual_end"]) if t["actual_end"] else "—"
            line = f"{owner}　·　完成 {ad}"
            text(s, content_x, info_y, content_w, info_h, line,
                 size=8, color=C["greenDeep"], valign="middle")

    # 底部 "另有 N 件" 提示
    if n_total > shown_n:
        text(s, x + 0.18, y + h - extra_h - 0.04, w - 0.36, extra_h,
             f"另有 {n_total - shown_n} 件,詳見 WBS 明細",
             size=8, italic=True, color=C["light"], align="right",
             valign="middle")


# ══════════════════════════════════════════════════════
# Slide N: 跨專案裁示事項彙整
# ══════════════════════════════════════════════════════
def slide_decisions(pres, r):
    s = _blank_slide(pres)
    high_n = sum(1 for d in r["decisions"] if d["severity"] == "高")
    mid_n  = sum(1 for d in r["decisions"] if d["severity"] == "中")
    page_header(s, "需主管裁示事項",
                f"高嚴重度 {high_n}　·　中嚴重度 {mid_n}　·　{fmt_date(r['report_date'])}",
                accent=C["red"])

    if not r["decisions"]:
        round_rect(s, 2.5, 2.2, 5.0, 1.2, C["greenSoft"],
                   line=C["divider"], line_w=0.5, radius=0.1)
        text(s, 2.5, 2.2, 5.0, 1.2, "✓ 目前無需裁示事項",
             size=18, bold=True, color=C["green"],
             align="center", valign="middle")
        return

    HX, HY, HW = 0.3, 0.92, 9.4
    COLW = [0.55, 1.55, 2.65, 2.30, 2.35]
    HDRS = ["嚴重度", "專案 / PM", "議題", "影響說明", "建議行動"]

    hx = HX
    for hd, w in zip(HDRS, COLW):
        rect(s, hx, HY, w, 0.34, C["bg"], line=C["divider"], line_w=0.4)
        text(s, hx, HY, w, 0.34, hd,
             size=9.5, bold=True, color=C["muted"],
             align="center", valign="middle")
        hx += w

    show = [d for d in r["decisions"] if d["severity"] in ("高", "中")][:7]
    avail = SLIDE_H - HY - 0.34 - 0.45
    RH = min(0.78, max(0.52, avail / max(len(show), 1)))

    for i, d in enumerate(show):
        y = HY + 0.34 + i * RH
        bg = C["white"] if i % 2 == 0 else C["rowAlt"]
        rect(s, HX, y, HW, RH, bg, line=C["divider"], line_w=0.3)

        sc = C["red"] if d["severity"] == "高" else C["amber"]
        dc = HX

        round_rect(s, dc + 0.10, y + 0.16, 0.36, RH - 0.32, sc, radius=0.20)
        text(s, dc + 0.10, y + 0.16, 0.36, RH - 0.32, d["severity"],
             size=10, bold=True, color=C["white"],
             align="center", valign="middle")
        dc += COLW[0]

        text(s, dc + 0.08, y + 0.07, COLW[1] - 0.16, 0.30, d["project"],
             size=10, bold=True, color=C["text"], valign="middle")
        text(s, dc + 0.08, y + 0.40, COLW[1] - 0.16, RH - 0.5,
             d["owner"], size=8.5, color=C["muted"])
        dc += COLW[1]

        text(s, dc + 0.10, y + 0.06, COLW[2] - 0.20, RH - 0.12,
             d["issue"],
             size=9.5, color=C["text"], valign="middle")
        dc += COLW[2]

        text(s, dc + 0.10, y + 0.06, COLW[3] - 0.20, RH - 0.12,
             d["impact"],
             size=8.5, italic=True, color=C["muted"], valign="middle")
        dc += COLW[3]

        text(s, dc + 0.10, y + 0.06, COLW[4] - 0.20, RH - 0.12,
             d["ask"],
             size=9, bold=True, color=sc, valign="middle")

    low_n = sum(1 for d in r["decisions"] if d["severity"] == "低")
    extra = (high_n + mid_n) - len(show)
    notes = []
    if extra > 0:
        notes.append(f"另有 {extra} 件高/中嚴重度事項")
    if low_n > 0:
        notes.append(f"另有 {low_n} 件低嚴重度事項")
    if notes:
        text(s, HX, HY + 0.34 + len(show) * RH + 0.10, HW, 0.22,
             "　·　".join(notes) + ",請參閱各專案 Page 2",
             size=9, color=C["light"], italic=True)


# ══════════════════════════════════════════════════════
# 結尾頁
# ══════════════════════════════════════════════════════
def slide_closing(pres, r):
    s = _blank_slide(pres)
    main_c = HEALTH[r["overall"]]["main"]

    rect(s, 0, 0, 0.10, SLIDE_H, main_c)
    rect(s, SLIDE_W - 0.06, 0, 0.06, SLIDE_H, C["divider"])

    text(s, 0, 1.9, SLIDE_W, 0.9, "謝謝聆聽",
         size=44, bold=True, color=C["text"], align="center")
    rect(s, 3.5, 3.0, 3.0, 0.008, C["divider"])
    text(s, 0, 3.20, SLIDE_W, 0.36,
         f"本期報告　{fmt_date(r['report_date'])}　　·　　下次回報　{fmt_date(r['next_report'])}",
         size=12, color=C["muted"], align="center")


# ══════════════════════════════════════════════════════
# 主入口
# ══════════════════════════════════════════════════════
def build_ppt(report: dict, output_path: str) -> None:
    pres = Presentation()
    pres.slide_width  = Inches(SLIDE_W)
    pres.slide_height = Inches(SLIDE_H)

    slide_cover(pres, report)
    slide_executive_summary(pres, report)
    for p in report["projects"]:
        slide_project_overview(pres, p, report)
        slide_project_tasks(pres, p, report)
    slide_closing(pres, report)

    pres.save(output_path)