"""
ppt_components.py — 共用樣式 + 基礎繪圖元件
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── 色盤 ────────────────────────────────────────────
def hex2rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

C = {k: hex2rgb(v) for k, v in {
    "white":     "FFFFFF",
    "bg":        "F8FAFC",
    "rowAlt":    "FAFAFA",
    "text":      "18181B",
    "muted":     "71717A",
    "light":     "A1A1AA",
    "divider":   "E4E4E7",
    "red":       "DC2626",
    "redSoft":   "FEF2F2",
    "redDeep":   "991B1B",
    "amber":     "D97706",
    "amberSoft": "FFFBEB",
    "amberDeep": "92400E",
    "green":     "16A34A",
    "greenSoft": "F0FDF4",
    "greenDeep": "14532D",
    "blue":      "2563EB",
    "blueSoft":  "EFF6FF",
}.items()}

HEALTH = {"RED":   {"main": C["red"],   "soft": C["redSoft"],   "deep": C["redDeep"]},
          "AMBER": {"main": C["amber"], "soft": C["amberSoft"], "deep": C["amberDeep"]},
          "GREEN": {"main": C["green"], "soft": C["greenSoft"], "deep": C["greenDeep"]}}

STATUS_COLOR = {
    "完成":   C["green"],
    "進行中": C["blue"],
    "延遲":   C["red"],
    "風險":   C["amber"],
    "待開始": C["light"],
}

FONT = "Calibri"

# ── 投影片尺寸 (16:9) ──────────────────────────────────
SLIDE_W = 10.0
SLIDE_H = 5.625


# ──────────────────────────────────────────────────────
# 基礎繪圖元件
# ──────────────────────────────────────────────────────
def _no_line(shape):
    shape.line.fill.background()


def _no_shadow(shape):
    shape.shadow.inherit = False


def rect(slide, x, y, w, h, fill, *, line=None, line_w=None):
    """方形矩形"""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        _no_line(sh)
    else:
        sh.line.color.rgb = line
        if line_w:
            sh.line.width = Pt(line_w)
    _no_shadow(sh)
    return sh


def round_rect(slide, x, y, w, h, fill, *, line=None, line_w=None, radius=0.06):
    """圓角矩形"""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        _no_line(sh)
    else:
        sh.line.color.rgb = line
        if line_w:
            sh.line.width = Pt(line_w)
    _no_shadow(sh)
    return sh


def circle(slide, cx, cy, d, fill, *, line=None):
    """以中心點 + 直徑畫圓"""
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(cx - d/2), Inches(cy - d/2),
                                 Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        _no_line(sh)
    else:
        sh.line.color.rgb = line
    _no_shadow(sh)
    return sh


def text(slide, x, y, w, h, content, *,
         size=10, bold=False, italic=False, color=None,
         align="left", valign="middle", font=FONT, line_space=1.0):
    """加文字框；自動處理對齊、垂直對齊、字體"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(36000)
    tf.margin_top  = tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    tf.vertical_anchor = {
        "top":    MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[valign]

    p = tf.paragraphs[0]
    p.alignment = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
    }[align]
    p.line_spacing = line_space

    r = p.add_run()
    r.text = content
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color or C["text"]
    return box


def divider(slide, x, y, w, *, color=None):
    rect(slide, x, y, w, 0.006, color or C["divider"])


# ──────────────────────────────────────────────────────
# 通用版型元件
# ──────────────────────────────────────────────────────
def page_header(slide, title, subtitle="", page_tag="", accent=None):
    """頂部色條 + 標題列"""
    accent = accent or C["text"]
    rect(slide, 0, 0, SLIDE_W, 0.06, accent)
    text(slide, 0.35, 0.10, 7.0, 0.55, title,
         size=20, bold=True, color=C["text"], valign="middle")
    if subtitle:
        text(slide, 6.6, 0.10, 3.1, 0.55, subtitle,
             size=9.5, color=C["muted"], align="right", valign="middle")
    rect(slide, 0, 0.72, SLIDE_W, 0.006, C["divider"])
    if page_tag:
        text(slide, 0, SLIDE_H - 0.30, SLIDE_W - 0.20, 0.20,
             page_tag, size=8, color=C["light"], align="right")


def kpi_card(slide, x, y, w, h, value, label, *,
             value_color=None, accent_color=None,
             value_size=24, label_size=9):
    """KPI 小卡片：頂部色條 + 大數字 + 標籤"""
    rect(slide, x, y, w, h, C["white"], line=C["divider"], line_w=0.5)
    if accent_color is not None:
        rect(slide, x, y, w, 0.06, accent_color)
    text(slide, x, y + 0.08, w, h * 0.55, str(value),
         size=value_size, bold=True, color=value_color or C["text"],
         align="center", valign="middle")
    text(slide, x, y + h - 0.32, w, 0.26, label,
         size=label_size, color=C["muted"], align="center")


def progress_bar(slide, x, y, w, h, pct, color, *,
                 vs_pct=None, vs_color=None):
    """進度條；可選 vs_pct (時間軸刻度線)"""
    pct = max(0, min(100, pct or 0))
    # 軌道
    round_rect(slide, x, y, w, h, C["divider"], radius=0.5)
    # 完成填色
    if pct > 0.5:
        fill_w = max(h, w * pct / 100)
        round_rect(slide, x, y, fill_w, h, color, radius=0.5)
    # 時間軸刻度（vs 線）
    if vs_pct is not None and 0 <= vs_pct <= 100:
        vx = x + w * vs_pct / 100
        # 一條從上突出的小刻度線
        rect(slide, vx - 0.015, y - 0.06, 0.030, h + 0.12,
             vs_color or C["text"])


def status_pill(slide, x, y, w, h, status):
    """狀態膠囊"""
    color = STATUS_COLOR.get(status, C["light"])
    round_rect(slide, x, y, w, h, color, radius=0.5)
    text(slide, x, y, w, h, status,
         size=7, bold=True, color=C["white"], align="center", valign="middle")


def health_badge(slide, cx, cy, d, level, *, with_ring=True):
    """大型健康度燈號（圓 + 可選外環）"""
    main = HEALTH[level]["main"]
    soft = HEALTH[level]["soft"]
    if with_ring:
        circle(slide, cx, cy, d, soft)
        circle(slide, cx, cy, d * 0.62, main)
    else:
        circle(slide, cx, cy, d, main)


# ──────────────────────────────────────────────────────
# 日期工具
# ──────────────────────────────────────────────────────
def fmt_date(d, fmt="%Y/%m/%d"):
    if d is None:
        return "—"
    return d.strftime(fmt)


def fmt_date_short(d):
    return fmt_date(d, "%m/%d")